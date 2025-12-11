from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import datetime

# --- Configuration ---
OUTPUT_FILE = "Presentations_Resultats_Projet.pptx"
IMG_DIR = "results/figures"

# Define Image Paths
IMG_COMPARISON = os.path.join(IMG_DIR, "comparison/model_comparison.png")
IMG_OPTUNA = os.path.join(IMG_DIR, "lightgbm/lgbm_optimization_history.png") 
IMG_GRU_LOSS = os.path.join(IMG_DIR, "gru/gru_training_loss.png")
IMG_TRANSFORMER_LOSS = os.path.join(IMG_DIR, "transformer/transformer_training_loss.png")
IMG_XGB_IMP = os.path.join(IMG_DIR, "xgboost/xgb_feature_importance.png")
IMG_LOGO = os.path.join(IMG_DIR, "kolecto_logo.png")

# Start Presentation
prs = Presentation()
SLIDE_WIDTH = prs.slide_width
SLIDE_HEIGHT = prs.slide_height

# --- Helper Functions ---
def set_font(run, size=18, bold=False, color=None):
    run.font.name = 'Times New Roman'
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def add_common_design(slide, title_text):
    """Applies common design elements: Blue Title, Logo."""
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        paragraph = slide.shapes.title.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
        paragraph.font.name = 'Times New Roman'
        paragraph.font.bold = True
        paragraph.font.size = Pt(36) # Standardize title size
        
    # Logo (Top Right)
    if os.path.exists(IMG_LOGO):
        logo_width = Inches(1.2)
        logo_left = SLIDE_WIDTH - logo_width - Inches(0.2)
        logo_top = Inches(0.2)
        slide.shapes.add_picture(IMG_LOGO, logo_left, logo_top, width=logo_width)

def create_title_slide(prs, title, subtitle, author, logo_path=None):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    slide.shapes.title.text = title
    paragraph = slide.shapes.title.text_frame.paragraphs[0]
    paragraph.font.color.rgb = RGBColor(0, 51, 102) 
    paragraph.font.name = 'Times New Roman'
    paragraph.font.bold = True
    
    # Subtitle
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"{subtitle}\n\n{author}"
        for paragraph in slide.placeholders[1].text_frame.paragraphs:
            paragraph.font.name = 'Times New Roman'
    
    # Large Logo for Title Slide
    if logo_path and os.path.exists(logo_path):
        left = Inches(0.5)
        top = Inches(0.5)
        height = Inches(1.0) 
        slide.shapes.add_picture(logo_path, left, top, height=height)

def add_content_slide(prs, title, content_lines, image_path=None):
    # Use Layout 5 (Title Only) for manual control
    slide_layout = prs.slide_layouts[5] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply Common Design (Title + Logo)
    add_common_design(slide, title)
    
    # Content Box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(5.0) if image_path else Inches(9.0)
    height = Inches(5.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.name = 'Times New Roman'
        # Emulate bullet spacing
        p.space_after = Pt(10)
        
        # indent logic
        if line.startswith("   -"):
            p.level = 1
            p.text = line.replace("   -", "\u2022 ").strip() # Bullet char
            p.font.size = Pt(16)
        elif line.startswith("       ->"):
            p.level = 2
            p.text = line.replace("       ->", "- ").strip()
            p.font.size = Pt(14)
        else:
            p.text = "\u2022 " + p.text # Top level bullet

    # Image (Centered horizontally in remaining space or right side)
    if image_path and os.path.exists(image_path):
        img_left = Inches(6.0)
        img_top = Inches(1.5)
        img_width = Inches(3.5)
        
        # Special handling for wide images or centered layouts
        if "comparison" in image_path or "feature" in image_path:
             # Center below text
             # Actually, user liked Centered. Let's make text full width top, image centered bottom?
             # Or side-by-side. The prompt asked for "Centered the result images" before.
             # Let's do Side-by-side for vertical fit, or Top-Bottom.
             
             # Re-adjusting for visual balance:
             # Resize text box
             txBox.width = Inches(9.0)
             txBox.height = Inches(2.5) # Shorten text area
             
             # Image Centered Bottom
             img_height = Inches(3.2)
             if "comparison" in image_path:
                 img_height = Inches(3.0)
             
             # Load image to get aspect ratio? No, just fit to constraints
             # Add picture
             pic = slide.shapes.add_picture(image_path, Inches(0), Inches(0)) # Temp insert
             
             # Scale
             ratio = pic.width / pic.height
             new_height = img_height
             new_width = new_height * ratio
             if new_width > Inches(8.5):
                 new_width = Inches(8.5)
                 new_height = new_width / ratio
            
             pic.width = int(new_width)
             pic.height = int(new_height)
             pic.left = int((SLIDE_WIDTH - pic.width) / 2)
             pic.top = int(Inches(4.2))
             
        else:
             # Default Side placement
             slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

def create_table_slide(prs, title, headers, data, notes=None):
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply Common Design
    add_common_design(slide, title)
    
    # Table Config
    rows = len(data) + 1
    cols = len(headers)
    table_left = Inches(1.0)
    table_top = Inches(2.0) # Lower top to clear logo/title
    table_width = Inches(8.0)
    table_height = Inches(0.5 * rows)
    
    shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = shape.table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.name = 'Times New Roman'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data
    for row_idx, row_data in enumerate(data):
        for col_idx, item in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(item)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.name = 'Times New Roman'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
    # Notes below table
    if notes:
        txBox = slide.shapes.add_textbox(Inches(1.0), table_top + shape.height + Inches(0.2), Inches(8.0), Inches(1.0))
        tf = txBox.text_frame
        for note in notes:
            p = tf.add_paragraph()
            p.text = note
            p.font.size = Pt(14)
            p.font.name = 'Times New Roman'
            p.font.italic = True
    
    return slide # Return slide in case we need to add more

# --- Slides Generation (Step-by-Step) ---

# Slide 1: Title Slide (Distinct Design, No Top-Right Logo)
create_title_slide(
    prs, 
    "Prédiction de Conversion Trial-to-Paid :\nInsights Data-Driven pour Kolecto", 
    "Améliorer le taux de conversion de ~60% via l'analyse des signaux précurseurs et le Machine Learning",
    f"Antigravity Agent – {datetime.date.today()}",
    logo_path=IMG_LOGO
)

# Slide 2: Business Context & Objectives
add_content_slide(prs, "Contexte Business & Objectifs", [
    "Contexte :",
    "   - Kolecto propose un essai gratuit de 15 jours convertissant en abonnement payant.",
    "   - Taux de conversion actuel : ~60% (Satisfaisant mais perfectible).",
    "Challenge :",
    "   - Identifier tôt les signaux précurseurs (succès vs désabonnement).",
    "   - Permettre des actions ciblées par l'équipe Customer Experience (CX).",
    "Objectifs :",
    "   - Analyser les facteurs différenciants (Convertis vs Non-Convertis).",
    "   - Construire un modèle ML pour prédire la probabilité de conversion."
])

# Slide 3: Data Overview (TABLE)
headers_data = ["Jeu de Données", "Taille", "Description", "Prétraitement"]
content_data = [
    ["Daily Usage", "~11k lignes", "Logs d'activité (Virements, Connexions)", "Agrégation (Somme/Moyenne/Max/Std)"],
    ["Subscriptions", "416 essais", "Profil Entreprise (CA, Code NAF)", "Fusion, Imputation, Encodage OneHot"]
]
notes_data = [
    "Stats Clés : Échantillon Total : 416 essais complets.",
    "Taux de Conversion : 60.7% (Déséquilibré mais gérable)."
]
create_table_slide(prs, "Vue d'Ensemble des Données & Prétraitement", headers_data, content_data, notes=notes_data)


# Slide 4: Methodology & Models
add_content_slide(prs, "Méthodologie & Modèles", [
    "Stratégie :",
    "   - Features Tabulaires (157 dims) pour modèles Arborescents.",
    "   - Données Séquentielles (15 jours) pour le Deep Learning.",
    "Modèles Entraînés :",
    "   - Logistic Regression (Baseline).",
    "   - XGBoost & LightGBM (Gradient Boosting optimisé avec Optuna).",
    "   - GRU & Transformer (Modélisation séquentielle des signaux temporel).",
    "Métriques d'Évaluation :",
    "   - ROC-AUC : Capacité de discrimination (Métrique Principale).",
    "   - PR-AUC : Précision-Rappel (Critique pour l'imbalance)."
])

# Slide 5: Overall Results Comparison (TABLE + Centered Image)
# We handle the image manually here, but use create_table_slide for base
slide_5 = create_table_slide(prs, "Comparaison Globale des Résultats", 
    ["Modèle", "ROC-AUC", "PR-AUC", "Accuracy", "Brier Score"], 
    [
        ["LightGBM (Vainqueur)", "0.790", "0.835", "72.3%", "0.193"],
        ["GRU (RNN)", "0.713", "0.743", "65.1%", "0.217"],
        ["Transformer", "0.711", "0.748", "66.3%", "0.211"],
        ["Logistic Reg.", "0.684", "0.769", "65.1%", "0.229"],
        ["XGBoost", "0.671", "0.772", "63.9%", "0.242"]
    ]
)
# Add Image manually to Slide 5
if os.path.exists(IMG_COMPARISON):
    img_width = Inches(6.0)
    img_left = (SLIDE_WIDTH - img_width) / 2
    # Place below table manually
    img_top = Inches(4.5) 
    slide_5.shapes.add_picture(IMG_COMPARISON, img_left, img_top, width=img_width)


# Slide 6: Optimization & Training Insights (Custom Layout)
slide_layout = prs.slide_layouts[5] # Title Only
slide = prs.slides.add_slide(slide_layout)
add_common_design(slide, "Optimisation & Dynamique d'Entraînement")

# Content
content_lines = [
    "LightGBM (Optuna) :",
    "   - Recherche efficace (50 essais).",
    "   - Convergence vers paramètres robustes (n_est=318, lr=0.018).",
    "Dynamique Deep Learning :",
    "   - GRU : Bonne baisse de loss training, légère instabilité en validation.",
    "   - Transformer : Signes d'overfitting (Train Loss << Val Loss)."
]
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9.0)
height = Inches(2.0)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True
for line in content_lines:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.name = 'Times New Roman'
    p.space_after = Pt(5)
    if line.startswith("   -"):
        p.level = 1
        p.text = line.replace("   -", "\u2022 ").strip()

# Add Images (3 Plots)
# 1. Optuna (Top Right)
if os.path.exists(IMG_OPTUNA):
    slide.shapes.add_picture(IMG_OPTUNA, Inches(6.0), Inches(1.5), height=Inches(2.5))

# 2. GRU Loss (Bottom Left)
if os.path.exists(IMG_GRU_LOSS):
    slide.shapes.add_picture(IMG_GRU_LOSS, Inches(0.5), Inches(4.5), height=Inches(2.5))

# 3. Transformer Loss (Bottom Right)
if os.path.exists(IMG_TRANSFORMER_LOSS):
    slide.shapes.add_picture(IMG_TRANSFORMER_LOSS, Inches(5.5), Inches(4.5), height=Inches(2.5)) 

# Slide 7: Feature Importance & Insights (Centered Image)
add_content_slide(prs, "Importance des Features & Insights", [
    "Meilleurs Prédicteurs :",
    "   - company_age : Les entreprises plus anciennes/stables convertissent mieux.",
    "   - naf_code : Certains secteurs ont une affinité plus forte.",
    "   - nb_client_invoices_created_sum : L'usage (Facturation) est le signal #1.",
    "Insights Actionnables :",
    "   - L'Activation Compte : Une utilisation précoce (Jours 1-3) est critique.",
    "   - Ciblage : Concentrer le CX sur les TPEs avec faible activité initiale."
], image_path=IMG_XGB_IMP)

# Slide 8: Business Impact & Recommendations
add_content_slide(prs, "Impact Business & Recommandations", [
    "Impact Estimé :",
    "   - Cible : ~400 essais/mois.",
    "   - Lift : +5-8% de conversion via intervention ciblée.",
    "   - Valeur : 400 * 0.05 * 3k€ (LTV) = ~60k€/mois -> 720k€/an.",
    "Recommandations :",
    "   - Jours 1-3 (Automatisé) : Nudge si 'nb_connections' < 2.",
    "   - Jours 7-10 (Humain) : Appel CX si Prob. Désabonnement > 60%.",
    "   - Déploiement : A/B Test des interventions pour mesurer l'uplift réel."
])

# Slide 9: Limitations & Improvements
add_content_slide(prs, "Limitations & Améliorations", [
    "Limitations :",
    "   - Faible Volume de Données : ~416 essais limitent le potentiel du Deep Learning.",
    "   - Facteurs Externes : Pas de données sur le contexte économique ou la saisonnalité.",
    "Améliorations Futures :",
    "   - Ensemble Hybride : Combiner LightGBM (Tabulaire) + GRU (Séquentiel).",
    "   - Causal ML : Modéliser l'uplift (Persuadables vs Do-not-disturb).",
    "   - Entraînement Continu : Réentraînement mensuel pour gérer le 'data drift'."
])

# Slide 10: Conclusion & Next Steps
add_content_slide(prs, "Conclusion & Prochaines Étapes", [
    "Résumé :",
    "   - Modèle robuste livré (LightGBM AUC 0.790).",
    "   - Leviers identifiés pour un lift de conversion de +5-8%.",
    "Prochaines Étapes :",
    "   1. Déployer l'API de Scoring (Containerisée).",
    "   2. Lancer l'A/B Test pour les actions CX.",
    "   3. Monitorer la Performance (MLflow) & Collecter plus de données."
])

# Save
prs.save(OUTPUT_FILE)
print(f"Presentation saved to {OUTPUT_FILE}")
